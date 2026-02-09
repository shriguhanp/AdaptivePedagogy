"""
Document Parser Service
=======================

Parses DOCX, PPTX, and PDF files for flashcard and quiz generation.
Uses simple, reliable libraries for fast text extraction.
"""

import os
import tempfile
from pathlib import Path
from typing import Optional

from src.logging import get_logger

logger = get_logger("DocumentParser")


def parse_docx(file_path: str) -> str:
    """Parse DOCX file and extract text."""
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
        return "\n".join(paragraphs)
    except ImportError:
        # Fallback to direct read if python-docx not available
        import zipfile
        from xml.etree import ElementTree as ET
        
        text_parts = []
        with zipfile.ZipFile(file_path) as z:
            for name in z.namelist():
                if name == "word/document.xml":
                    tree = ET.fromstring(z.read(name))
                    # Remove namespace
                    for elem in tree.iter():
                        if '}' in elem.tag:
                            elem.tag = elem.tag.split('}', 1)[1]
                    for para in tree.findall('.//p'):
                        text = "".join(para.itertext()).strip()
                        if text:
                            text_parts.append(text)
        return "\n".join(text_parts)


def parse_pptx(file_path: str) -> str:
    """Parse PPTX file and extract text from slides."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
        return "\n".join(text_parts)
    except ImportError:
        # Fallback for older python-pptx versions
        import zipfile
        from xml.etree import ElementTree as ET
        
        text_parts = []
        with zipfile.ZipFile(file_path) as z:
            for name in z.namelist():
                if name.startswith("ppt/slides/slide"):
                    try:
                        tree = ET.fromstring(z.read(name))
                        for elem in tree.iter():
                            if '}' in elem.tag:
                                elem.tag = elem.tag.split('}', 1)[1]
                        text = "".join(elem.itertext() for elem in tree.findall('.//a:t')).strip()
                        if text:
                            text_parts.append(text)
                    except:
                        pass
        return "\n".join(text_parts)


def parse_pdf(file_path: str) -> str:
    """Parse PDF file and extract text."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        text_parts = []
        for page in doc:
            text = page.get_text()
            if text.strip():
                text_parts.append(text.strip())
        doc.close()
        return "\n".join(text_parts)
    except ImportError:
        logger.warning("PyMuPDF not available, trying alternative PDF parsing")
        # Fallback to basic PDF parsing
        try:
            from pdfminer.high_level import extract_text
            return extract_text(file_path)
        except:
            return ""


def parse_document(file_path: str, max_chars: int = 5000) -> str:
    """
    Parse a document file and extract its text content.
    
    Args:
        file_path: Path to the document file
        max_chars: Maximum characters to return (for speed)
        
    Returns:
        Extracted text content
    """
    path = Path(file_path)
    
    if not path.exists():
        raise FileNotFoundError(f"Document not found: {file_path}")
    
    ext = path.suffix.lower()
    text = ""
    
    try:
        if ext in [".docx"]:
            text = parse_docx(str(path))
        elif ext in [".pptx", ".ppt"]:
            text = parse_pptx(str(path))
        elif ext in [".pdf"]:
            text = parse_pdf(str(path))
        else:
            raise ValueError(f"Unsupported file format: {ext}")
        
        # Truncate for speed if needed
        if len(text) > max_chars:
            logger.info(f"Truncating document from {len(text)} to {max_chars} chars")
            text = text[:max_chars]
        
        return text
        
    except Exception as e:
        logger.error(f"Failed to parse {file_path}: {e}")
        raise


async def parse_uploaded_document(file_content: bytes, filename: str, max_chars: int = 5000) -> str:
    """
    Parse an uploaded file from memory.
    
    Args:
        file_content: Raw file bytes
        filename: Original filename
        max_chars: Maximum characters to return
        
    Returns:
        Extracted text content
    """
    # Create temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=filename) as tmp:
        tmp.write(file_content)
        tmp_path = tmp.name
    
    try:
        return parse_document(tmp_path, max_chars)
    finally:
        # Clean up
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)


# Singleton parser
_parser: Optional["SimpleDocumentParser"] = None


class SimpleDocumentParser:
    """Simple document parser wrapper."""
    
    @staticmethod
    def parse(file_path: str, max_chars: int = 5000) -> str:
        return parse_document(file_path, max_chars)
    
    @staticmethod
    async def parse_async(file_content: bytes, filename: str, max_chars: int = 5000) -> str:
        return await parse_uploaded_document(file_content, filename, max_chars)
    
    # Alias for parse_async (for backward compatibility)
    @staticmethod
    async def parse_uploaded_file(file_content: bytes, filename: str, max_chars: int = 5000) -> str:
        return await parse_uploaded_document(file_content, filename, max_chars)


def get_document_parser() -> SimpleDocumentParser:
    """Get the document parser instance."""
    return SimpleDocumentParser()
